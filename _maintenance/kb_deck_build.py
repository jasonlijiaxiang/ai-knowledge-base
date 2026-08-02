#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""讲义生成器：把结构化内容 JSON 渲染成符合全库设计系统的讲义 PPTX（需 python-pptx）。

**为什么要有它**（2026-08-02 建两个新模块时立）：此前 19 册讲义都是逐页手写的，一册
六七十到一百页。手写在单册尺度上可行，到「同时建两册、每册十来章」就不可行了——不是写不动，
是**设计令牌会在几百页里慢慢漂**：字号差半档、深浅页归属记错、页脚章节 ID 忘了换、
封面章节条停在旧章数。这些正是 `audit_pptx.py` 十七项里反复出现的那几类。

所以把设计系统的**机械部分**编码进来，人只写内容：

  · 令牌（尺寸/字体/字号档/色板/边距）来自 `ppt-design-system.md` §1–§3，写死在本文件顶部；
  · 明暗节奏按 §2 自动归属——封面、章节过渡、动手做、总收束是深色页，其余浅色；
  · 页脚按 §7 自动写「模块 · 章节 ID」与页码，**页码即放映位**，天然满足 check_footer_pagenum；
  · `docProps/app.xml` 由 python-pptx 自己维护，避免 §7 那条「插页漏改 app.xml 导致
    PowerPoint 弹修复」的老账。

**它不替人做的事**：写什么、每章几页、版式选哪种、类比与案例好不好——那些是 ppt-rules
和 core-rules §六 的活。生成器只保证「长相合规」，`audit_pptx.py` 与渲染目检照跑不误。

内容 JSON 结构（一页一个对象，按放映序排）：

    {"module": "Predictive-AI-MLOps", "display": "AI 知识库",
     "pages": [{"kind": "cover"|"toc"|"chapter"|"goals"|"points"|"table"|
                        "steps"|"qa"|"handson"|"summary"|"sources"|"recap",
                "chapter_id": "pam-what-why", "title": "...", ...}]}

各 kind 需要的字段见对应 render 函数的文档串。

用法:
    python3 kb_deck_build.py <内容.json> <输出.pptx>

生成后必须接：`audit_pptx.py <输出.pptx>`（十七项）+ 逐页渲染目检——
本脚本只管令牌，管不了「文字撑破卡片」那一格（错题本里那条只有渲染看得见）。
"""
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---- 设计令牌（ppt-design-system §1）----
W, H = 13.333, 7.5
C = {
    "dark": "0B2540", "ink": "1A2733", "ink2": "51606E", "cyan": "128199",
    "cyan_lt": "1CA3B8", "cyan_dk": "0E4D64", "orange": "F2A65A", "green": "2F855A",
    "card": "F3F8FA", "card2": "EAF3F5", "rule": "D6E1E8", "faint": "9FB3C0",
    "white": "FFFFFF",
}
TITLE_FONT, BODY_FONT, CODE_FONT = "Cambria", "Calibri", "Courier New"

# 自读档（KB-CONFIG「阅读场景」=自读；ppt-design-system §3）
FS = {"cover": 44, "chapter": 36, "h1": 28, "sub": 20, "body": 13.5,
      "table": 12, "note": 11, "foot": 9, "code": 12}

MARGIN = 0.62          # 左右安全边距 ≥0.55
TOP = 0.5              # 标题区起点
FOOT_Y = 6.95          # 页脚区
BODY_TOP = 1.55        # 正文起点（标题区之下）


def rgb(k):
    return RGBColor.from_string(C[k])


def _tf(shape, size, color, bold=False, font=BODY_FONT, align=PP_ALIGN.LEFT,
        spacing=1.0):
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = rgb(color)
            r.font.bold = bold
            r.font.name = font
    return tf


def box(slide, x, y, w, h, text, size, color, **kw):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    sh.text_frame.text = text
    _tf(sh, size, color, **kw)
    return sh


def para(shape, text, size, color, bold=False, space_before=6, bullet=None):
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = ("%s %s" % (bullet, text)) if bullet else text
    p.space_before = Pt(space_before)
    p.line_spacing = 1.25
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = rgb(color)
        r.font.bold = bold
        r.font.name = BODY_FONT
    return p


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    if line:
        sh.line.color.rgb = rgb(line)
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


DARK_KINDS = {"cover", "chapter", "handson", "recap"}


def new_slide(prs, kind):
    s = prs.slides.add_slide(prs.slide_layouts[6])      # 纯空白版式
    if kind in DARK_KINDS:
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb("dark")
    return s


def _k():
    """当前画布相对标准画布（13.333in）的比例。

    本模块的版式常量按 13.333×7.5 写死，改 W/H 只改了一部分——`head` 的眉题宽 11.5、
    `footer` 的左栏宽 8.0 与页码栏宽 1.2 都是绝对英寸。往 10×5.625 的册插页时，
    这三处会顶到 12.0in，被 audit 的检查项 16 判成「大画布坐标页嫁接进小画布册」。
    这里统一按比例取，W=13.333 时恒为 1，既有产物字节不变。
    """
    return W / 13.333


def footer(slide, mod, chap, num, dark=False, foot=None):
    """ppt-design-system §7：左下页脚，右下页码；页码＝放映位。

    左下优先用 `foot`（形如「第 3 章 · 把业务问题切成可学的题 · 本章小结」）——
    **audit 靠这一行把学习目标/动手做/对练/小结归到章**，只写「模块 · 章节 ID」它归不上，
    会误报「每章固定元素缺失」。既有 19 册用的就是这个写法，照它来。
    """
    col = "faint" if dark else "ink2"
    left = foot or ("%s · %s" % (mod, chap) if chap else mod)
    k = _k()
    box(slide, MARGIN, FOOT_Y, 8.0 * k, 0.3 * k, left, FS["foot"], col)
    b = box(slide, W - MARGIN - 1.2 * k, FOOT_Y, 1.2 * k, 0.3 * k, "p.%d" % num,
            FS["foot"], col, align=PP_ALIGN.RIGHT)
    return b


def head(slide, title, dark=False, eyebrow=None):
    k = _k()
    if eyebrow:
        box(slide, MARGIN, TOP - 0.28 * k, 11.5 * k, 0.3 * k, eyebrow, FS["foot"] + 1.5,
            "cyan_lt" if dark else "cyan", bold=True)
    box(slide, MARGIN, TOP, W - 2 * MARGIN, 0.95 * k, title, FS["h1"],
        "white" if dark else "ink", bold=True, font=TITLE_FONT)


# ---------------- 各 kind ----------------

def r_cover(s, p, ctx):
    """封面契约（§4）：眉题 / 主标题 / 副题 / 章节概览 / 读者定位 + 核实日期。"""
    box(s, MARGIN, 0.75, 11.0, 0.35, "%s · 讲义式 PPT" % ctx["display"],
        13, "cyan_lt")
    box(s, MARGIN, 1.35, 11.6, 1.9, p["title"], FS["cover"], "white",
        bold=True, font=TITLE_FONT, spacing=1.08)
    box(s, MARGIN, 3.45, 11.6, 0.6, p["subtitle"], 21, "cyan_lt")
    y = 4.4
    groups = p.get("groups") or []
    if groups:
        gw = (W - 2 * MARGIN - 0.2 * (len(groups) - 1)) / len(groups)
        for i, g in enumerate(groups):
            x = MARGIN + i * (gw + 0.2)
            sh = rect(s, x, y, gw, 0.82, "cyan_dk")
            sh.text_frame.text = g["label"]
            _tf(sh, 12.5, "white", bold=True, align=PP_ALIGN.CENTER)
            para(sh, g["desc"], 10.5, "cyan_lt", space_before=2)
            for pa in sh.text_frame.paragraphs:
                pa.alignment = PP_ALIGN.CENTER
        y += 1.15
    box(s, MARGIN, y + 0.35, 11.6, 0.35, p["audience"], 11.5, "faint")
    box(s, MARGIN, y + 0.72, 11.6, 0.35, "核实日期 %s" % p["verified"], 11, "faint")


def r_toc(s, p, ctx):
    head(s, "导览")
    cols, items = 3, p["items"]
    per = (len(items) + cols - 1) // cols
    cw = (W - 2 * MARGIN - 0.3 * (cols - 1)) / cols
    for c in range(cols):
        chunk = items[c * per:(c + 1) * per]
        if not chunk:
            continue
        x = MARGIN + c * (cw + 0.3)
        sh = s.shapes.add_textbox(Inches(x), Inches(BODY_TOP), Inches(cw), Inches(4.9))
        sh.text_frame.text = ""
        for it in chunk:
            para(sh, "%s  %s" % (it["no"], it["title"]), 12.5, "ink",
                 bold=True, space_before=9)
            if it.get("gist"):
                para(sh, it["gist"], 10.5, "ink2", space_before=1)
    if p.get("note"):
        box(s, MARGIN, 6.45, 11.6, 0.4, p["note"], FS["note"], "ink2")


def r_chapter(s, p, ctx):
    box(s, MARGIN, 2.35, 11.6, 0.5, p["no"], 15, "cyan_lt", bold=True)
    box(s, MARGIN, 2.85, 11.6, 1.2, p["title"], FS["chapter"], "white",
        bold=True, font=TITLE_FONT)
    box(s, MARGIN, 4.35, 10.5, 1.0, p["question"], 17, "orange")


def r_goals(s, p, ctx):
    head(s, p["title"], eyebrow=p.get("eyebrow"))
    sh = s.shapes.add_textbox(Inches(MARGIN), Inches(BODY_TOP),
                              Inches(W - 2 * MARGIN), Inches(4.6))
    sh.text_frame.text = ""
    for g in p["goals"]:
        para(sh, g, 14.5, "ink", space_before=13, bullet="▸")


def r_points(s, p, ctx):
    """要点卡：2–4 张，每张标题 + 说明。避免四宫格定义卡片的滥用见 §5。"""
    head(s, p["title"], eyebrow=p.get("eyebrow"))
    if p.get("lead"):
        box(s, MARGIN, BODY_TOP - 0.28, W - 2 * MARGIN, 0.5, p["lead"],
            FS["sub"] - 3, "ink2")
    cards = p["cards"]
    n = len(cards)
    top = BODY_TOP + (0.5 if p.get("lead") else 0)
    cw = (W - 2 * MARGIN - 0.28 * (n - 1)) / n
    ch = 6.6 - top
    for i, c in enumerate(cards):
        x = MARGIN + i * (cw + 0.28)
        sh = rect(s, x, top, cw, ch, "card")
        sh.text_frame.margin_left = Inches(0.18)
        sh.text_frame.margin_right = Inches(0.18)
        sh.text_frame.margin_top = Inches(0.16)
        sh.text_frame.text = c["h"]
        _tf(sh, 14, "cyan_dk", bold=True)
        for line in c["lines"]:
            para(sh, line, FS["body"], "ink", space_before=7)
        sh.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def r_steps(s, p, ctx):
    """机制/流程：3–7 步，左到右。"""
    head(s, p["title"], eyebrow=p.get("eyebrow"))
    if p.get("lead"):
        box(s, MARGIN, BODY_TOP - 0.28, W - 2 * MARGIN, 0.5, p["lead"],
            FS["sub"] - 3, "ink2")
    steps = p["steps"]
    n = len(steps)
    top = BODY_TOP + (0.55 if p.get("lead") else 0.1)
    gap = 0.34
    sw = (W - 2 * MARGIN - gap * (n - 1)) / n
    for i, st in enumerate(steps):
        x = MARGIN + i * (sw + gap)
        sh = rect(s, x, top, sw, 2.5, "card2")
        sh.text_frame.margin_left = Inches(0.14)
        sh.text_frame.margin_right = Inches(0.14)
        sh.text_frame.text = "%d" % (i + 1)
        _tf(sh, 12, "cyan", bold=True)
        para(sh, st["h"], 13.5, "ink", bold=True, space_before=3)
        para(sh, st["d"], 11.5, "ink2", space_before=4)
        if i < n - 1:
            box(s, x + sw + 0.02, top + 1.0, gap - 0.04, 0.4, "→", 15, "cyan",
                align=PP_ALIGN.CENTER)
    if p.get("takeaway"):
        sh = rect(s, MARGIN, top + 2.75, W - 2 * MARGIN, 0.85, "card")
        sh.text_frame.margin_left = Inches(0.18)
        sh.text_frame.text = p["takeaway"]
        _tf(sh, 13, "cyan_dk", bold=True)


def r_table(s, p, ctx):
    head(s, p["title"], eyebrow=p.get("eyebrow"))
    if p.get("lead"):
        box(s, MARGIN, BODY_TOP - 0.28, W - 2 * MARGIN, 0.5, p["lead"],
            FS["sub"] - 3, "ink2")
    cols, rows = p["cols"], p["rows"]
    top = BODY_TOP + (0.5 if p.get("lead") else 0.05)
    gt = s.shapes.add_table(len(rows) + 1, len(cols), Inches(MARGIN), Inches(top),
                            Inches(W - 2 * MARGIN), Inches(0.4 * (len(rows) + 1)))
    t = gt.table
    if p.get("widths"):
        total = sum(p["widths"])
        for i, wgt in enumerate(p["widths"]):
            t.columns[i].width = Inches((W - 2 * MARGIN) * wgt / total)
    for j, c in enumerate(cols):
        cell = t.cell(0, j)
        cell.text = c
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb("cyan_dk")
        for pa in cell.text_frame.paragraphs:
            for r in pa.runs:
                r.font.size = Pt(FS["table"])
                r.font.bold = True
                r.font.color.rgb = rgb("white")
                r.font.name = BODY_FONT
    for i, row in enumerate(rows, 1):
        for j, v in enumerate(row):
            cell = t.cell(i, j)
            cell.text = v
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("white" if i % 2 else "card")
            for pa in cell.text_frame.paragraphs:
                for r in pa.runs:
                    r.font.size = Pt(FS["table"])
                    r.font.color.rgb = rgb("ink")
                    r.font.name = BODY_FONT
    if p.get("note"):
        box(s, MARGIN, 6.45, W - 2 * MARGIN, 0.4, p["note"], FS["note"], "ink2")


def r_qa(s, p, ctx):
    head(s, p["title"], eyebrow=p.get("eyebrow"))
    sh = s.shapes.add_textbox(Inches(MARGIN), Inches(BODY_TOP),
                              Inches(W - 2 * MARGIN), Inches(4.9))
    sh.text_frame.text = ""
    for qa in p["qa"]:
        para(sh, "Q  " + qa["q"], 13.5, "cyan_dk", bold=True, space_before=14)
        para(sh, "A  " + qa["a"], 12.5, "ink", space_before=3)


def r_handson(s, p, ctx):
    head(s, p["title"], dark=True, eyebrow=p.get("eyebrow"))
    sh = s.shapes.add_textbox(Inches(MARGIN), Inches(BODY_TOP),
                              Inches(W - 2 * MARGIN), Inches(4.2))
    sh.text_frame.text = ""
    for t in p["tasks"]:
        para(sh, t, 14, "white", space_before=13, bullet="▸")
    if p.get("evidence"):
        box(s, MARGIN, 6.05, W - 2 * MARGIN, 0.6,
            "成功证据　%s" % p["evidence"], 12.5, "orange")


def r_summary(s, p, ctx):
    head(s, p["title"], eyebrow=p.get("eyebrow"))
    sh = s.shapes.add_textbox(Inches(MARGIN), Inches(BODY_TOP),
                              Inches(W - 2 * MARGIN), Inches(3.9))
    sh.text_frame.text = ""
    for i, t in enumerate(p["takeaways"], 1):
        para(sh, "%d　%s" % (i, t), 14, "ink", space_before=14)
    if p.get("next"):
        sh2 = rect(s, MARGIN, 5.85, W - 2 * MARGIN, 0.8, "card")
        sh2.text_frame.margin_left = Inches(0.18)
        sh2.text_frame.text = "下一章　%s" % p["next"]
        _tf(sh2, 12.5, "cyan_dk", bold=True)


def r_recap(s, p, ctx):
    head(s, p["title"], dark=True)
    n = len(p["lines"])
    half = (n + 1) // 2
    for c, chunk in enumerate((p["lines"][:half], p["lines"][half:])):
        if not chunk:
            continue
        x = MARGIN + c * ((W - 2 * MARGIN) / 2 + 0.15)
        sh = s.shapes.add_textbox(Inches(x), Inches(BODY_TOP),
                                  Inches((W - 2 * MARGIN) / 2 - 0.15), Inches(4.9))
        sh.text_frame.text = ""
        for l in chunk:
            para(sh, l, 12.5, "white", space_before=10, bullet="·")


def r_sources(s, p, ctx):
    head(s, p["title"])
    if p.get("lead"):
        box(s, MARGIN, BODY_TOP - 0.28, W - 2 * MARGIN, 0.5, p["lead"],
            FS["note"] + 1, "ink2")
    sh = s.shapes.add_textbox(Inches(MARGIN), Inches(BODY_TOP + 0.35),
                              Inches(W - 2 * MARGIN), Inches(4.6))
    sh.text_frame.text = ""
    for g in p["groups"]:
        para(sh, g["h"], 12.5, "cyan_dk", bold=True, space_before=11)
        for it in g["items"]:
            para(sh, it, 10.5, "ink2", space_before=2)


RENDER = {"cover": r_cover, "toc": r_toc, "chapter": r_chapter, "goals": r_goals,
          "points": r_points, "steps": r_steps, "table": r_table, "qa": r_qa,
          "handson": r_handson, "summary": r_summary, "recap": r_recap,
          "sources": r_sources}


APP_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"><TotalTime>0</TotalTime><Application>python-pptx</Application><PresentationFormat>Widescreen</PresentationFormat><Slides>{n}</Slides><Notes>0</Notes><HiddenSlides>0</HiddenSlides><MMClips>0</MMClips><ScaleCrop>false</ScaleCrop><HeadingPairs><vt:vector size="4" baseType="variant"><vt:variant><vt:lpstr>Theme</vt:lpstr></vt:variant><vt:variant><vt:i4>1</vt:i4></vt:variant><vt:variant><vt:lpstr>Slide Titles</vt:lpstr></vt:variant><vt:variant><vt:i4>{n}</vt:i4></vt:variant></vt:vector></HeadingPairs><TitlesOfParts><vt:vector size="{n1}" baseType="lpstr"><vt:lpstr>Office Theme</vt:lpstr>{titles}</vt:vector></TitlesOfParts><LinksUpToDate>false</LinksUpToDate><SharedDoc>false</SharedDoc><HyperlinksChanged>false</HyperlinksChanged></Properties>"""


def _esc(t):
    return (t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def fix_app_xml(path, titles):
    """重建 docProps/app.xml。

    python-pptx 不维护它，Slides 会留在 0——**只有 PowerPoint 校验这份登记表**，
    渲染与常规脚本全部无视，登记与实际不符时 PowerPoint 弹「需要修复」
    （ppt-design-system §7 与 audit 检查项 9 记的就是这笔账）。
    """
    import shutil
    import zipfile
    n = len(titles)
    body = APP_XML.format(n=n, n1=n + 1,
                          titles="".join("<vt:lpstr>%s</vt:lpstr>" % _esc(t) for t in titles))
    tmp = path + ".tmp"
    zin = zipfile.ZipFile(path)
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for it in zin.infolist():
            data = zin.read(it.filename)
            if it.filename == "docProps/app.xml":
                data = body.encode("utf-8")
            zout.writestr(it, data)
    zin.close()
    shutil.move(tmp, path)


def build(spec, out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    ctx = {"display": spec.get("display", "AI 知识库"), "module": spec["module"]}
    for i, p in enumerate(spec["pages"], 1):
        kind = p["kind"]
        if kind not in RENDER:
            raise SystemExit("第 %d 页 kind「%s」不认识" % (i, kind))
        s = new_slide(prs, kind)
        RENDER[kind](s, p, ctx)
        if kind != "cover":
            footer(s, spec["module"], p.get("chapter_id", ""), i,
                   dark=kind in DARK_KINDS, foot=p.get("foot"))
    prs.save(out)
    fix_app_xml(out, [p.get("title") or p.get("no", "") for p in spec["pages"]])
    return len(spec["pages"])


def main(argv):
    if len(argv) < 3:
        print(__doc__)
        return 2
    spec = json.load(open(argv[1], encoding="utf-8"))
    n = build(spec, argv[2])
    print("已生成 %s：%d 页" % (argv[2], n))
    print("接着跑：python3 _maintenance/audit_pptx.py %s   然后逐页渲染目检" % argv[2])
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
