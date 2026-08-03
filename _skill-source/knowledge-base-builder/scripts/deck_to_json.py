#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""讲义 → 生成器源 JSON 反解器（需 `python-pptx`）。

**为什么要有它**（2026-08-03 工程审视时补）：`kb_deck_build.py` 立的规矩是「产物从源现场
生成」，与库里对技能包、网页产物一以贯之的纪律一致。但 2026-08-02 生成的两册
（AI-Governance 117 页、Predictive-AI-MLOps 121 页）的**整册内容 JSON 既没进版本库、
也不在磁盘上**——组装脚本读一个临时 JSON、写另一个临时 JSON，两个都没留。

结果是生成器唯一的卖点拿不到：设计系统一改，这两册和另外 19 册手写册处境完全相同，
只能做 zip 级手术。也就是说，**全库唯一能兑现「改令牌就重渲染」的地方，恰恰兑现不了**。

原件找不回来，但 `kb_deck_build.py` 的渲染是确定性的——每种 kind 的形状顺序、坐标、
填充色都写死在版式常量里。所以可以逆着走一遍：按坐标与填充色认出 kind，再把文字
抽回字段。**这不是猜**，判据全部来自 `kb_deck_build.py` 自己的常量。

验收不靠肉眼，靠往返：

    反解 → kb_deck_build.py 重渲染 → kb_deck_text.py 逐页文本比对 → 完全相等才算忠实

`--verify` 就是把这条链跑一遍（重渲染写临时文件，**原册字节一律不动**）。

用法:
    python3 deck_to_json.py <讲义.pptx> <输出.json>
    python3 deck_to_json.py <讲义.pptx> <输出.json> --verify
退出码: 0 = 反解成功（带 --verify 时往返也一致）, 1 = 有页认不出或往返不一致。
"""
import json
import os
import re
import subprocess
import sys
import tempfile

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))

# 与 kb_deck_build.py 的版式常量一一对应。改那边就要改这边——
# 两处漂了，反解出来的 JSON 会静默变形，只有 --verify 抓得住。
MARGIN, TOP, BODY_TOP, FOOT_Y = 0.62, 0.5, 1.55, 6.95
CARD, CARD2, CYAN_DK, DARK = "F3F8FA", "EAF3F5", "0E4D64", "0B2540"
EPS = 0.02          # 英寸容差（EMU 取整误差远小于此）


def at(v, target):
    return abs(v - target) < EPS


class Page:
    """一页的形状清单，按创建顺序，坐标换算成英寸。"""

    def __init__(self, slide, idx):
        self.idx = idx
        self.dark = False
        try:
            if slide.background.fill.type == 1:
                self.dark = str(slide.background.fill.fore_color.rgb) == DARK
        except Exception:
            pass
        self.boxes, self.rects, self.table = [], [], None
        for sh in slide.shapes:
            if sh.has_table:
                self.table = sh.table
                continue
            if not sh.has_text_frame:
                continue
            item = {
                "x": sh.left / 914400, "y": sh.top / 914400,
                "w": sh.width / 914400, "h": sh.height / 914400,
                "paras": [p.text for p in sh.text_frame.paragraphs],
                "bolds": [any(r.font.bold for r in p.runs) for p in sh.text_frame.paragraphs],
                "sizes": [(p.runs[0].font.size.pt if p.runs and p.runs[0].font.size else None)
                          for p in sh.text_frame.paragraphs],
            }
            item["text"] = "\n".join(item["paras"])
            fill = None
            try:
                if sh.fill.type == 1:
                    fill = str(sh.fill.fore_color.rgb)
            except Exception:
                pass
            (self.rects if fill else self.boxes).append(dict(item, fill=fill))

    # ---- 定位小工具 ----
    def find(self, **kw):
        for b in self.boxes:
            if all(at(b[k], v) for k, v in kw.items()):
                return b
        return None

    def all(self, **kw):
        return [b for b in self.boxes
                if all(at(b[k], v) for k, v in kw.items())]

    def rects_fill(self, fill):
        return [r for r in self.rects if r["fill"] == fill]


def strip_bullet(t, bullet):
    return t[len(bullet) + 1:] if t.startswith(bullet + " ") else t


def parse(page, module):
    """返回 (页对象, 页脚左栏文字)。认不出就抛异常——静默猜错比报错危险。"""
    p = {}
    foot_left = page.find(x=MARGIN, y=FOOT_Y, w=8.0)
    foot = foot_left["text"] if foot_left else None

    # ---- cover：唯一没有页脚的一页 ----
    if foot is None:
        head = page.find(x=MARGIN, y=0.75, w=11.0)
        ver = [b for b in page.boxes if b["text"].startswith("核实日期 ")]
        if not head or not ver:
            raise ValueError("第 %d 页没有页脚，也不像封面" % page.idx)
        title = page.find(x=MARGIN, y=1.35, w=11.6)
        sub = page.find(x=MARGIN, y=3.45, w=11.6)
        groups = [{"label": r["paras"][0], "desc": r["paras"][1]}
                  for r in page.rects_fill(CYAN_DK) if at(r["y"], 4.4)]
        # audience 在 verified 之上 0.37in
        aud = [b for b in page.boxes
               if at(b["y"], ver[0]["y"] - 0.37) and at(b["x"], MARGIN)]
        p = {"kind": "cover", "title": title["text"], "subtitle": sub["text"],
             "audience": aud[0]["text"] if aud else "",
             "verified": ver[0]["text"][len("核实日期 "):]}
        if groups:
            p["groups"] = groups
        p["_display"] = head["text"].rsplit(" · ", 1)[0]
        return p, None

    title_box = page.find(x=MARGIN, y=TOP, w=13.333 - 2 * MARGIN)
    eyebrow = page.find(x=MARGIN, y=TOP - 0.28, w=11.5)
    lead = page.find(x=MARGIN, y=BODY_TOP - 0.28, w=13.333 - 2 * MARGIN)

    def base(kind):
        d = {"kind": kind, "title": title_box["text"]}
        if eyebrow:
            d["eyebrow"] = eyebrow["text"]
        return d

    # ---- chapter：深色 + 章号在 y=2.35 ----
    no = page.find(x=MARGIN, y=2.35, w=11.6)
    if page.dark and no:
        t = page.find(x=MARGIN, y=2.85, w=11.6)
        q = page.find(x=MARGIN, y=4.35, w=10.5)
        return {"kind": "chapter", "no": no["text"], "title": t["text"],
                "question": q["text"]}, foot

    if title_box is None:
        raise ValueError("第 %d 页没有标题区，无法归类" % page.idx)

    # ---- table：唯一带 GraphicFrame 表格的 ----
    if page.table is not None:
        t = page.table
        d = base("table")
        d["cols"] = [t.cell(0, j).text for j in range(len(t.columns))]
        d["rows"] = [[t.cell(i, j).text for j in range(len(t.columns))]
                     for i in range(1, len(t.rows))]
        widths = [c.width for c in t.columns]
        if len(set(widths)) > 1:
            d["widths"] = [round(w / 914400, 4) for w in widths]
        if lead:
            d["lead"] = lead["text"]
        note = page.find(x=MARGIN, y=6.45)   # 表格页 12.09 宽、导览页 11.6 宽，只认位置
        if note:
            d["note"] = note["text"]
        return d, foot

    # ---- steps：card2 填充的步骤盒 ----
    if page.rects_fill(CARD2):
        d = base("steps")
        if lead:
            d["lead"] = lead["text"]
        boxes = sorted(page.rects_fill(CARD2), key=lambda r: r["x"])
        d["steps"] = [{"h": r["paras"][1], "d": r["paras"][2]} for r in boxes]
        take = [r for r in page.rects_fill(CARD) if r["w"] > 10]
        if take:
            d["takeaway"] = take[0]["text"]
        return d, foot

    # ---- points：card 填充的要点卡（占满正文高度）----
    cards = [r for r in page.rects_fill(CARD) if r["h"] > 2.5]
    if cards:
        d = base("points")
        if lead:
            d["lead"] = lead["text"]
        d["cards"] = [{"h": r["paras"][0], "lines": r["paras"][1:]}
                      for r in sorted(cards, key=lambda r: r["x"])]
        return d, foot

    # ---- recap：深色 + 两栏 ----
    if page.dark:
        cols = sorted(page.all(y=BODY_TOP, w=(13.333 - 2 * MARGIN) / 2 - 0.15),
                      key=lambda b: b["x"])
        if cols:
            lines = []
            for c in cols:
                lines += [strip_bullet(t, "·") for t in c["paras"] if t]
            return {"kind": "recap", "title": title_box["text"], "lines": lines}, foot
        body = page.find(x=MARGIN, y=BODY_TOP, w=13.333 - 2 * MARGIN)
        if body:
            d = base("handson")
            d["tasks"] = [strip_bullet(t, "▸") for t in body["paras"] if t]
            ev = [b for b in page.boxes if b["text"].startswith("成功证据　")]
            if ev:
                d["evidence"] = ev[0]["text"][len("成功证据　"):]
            return d, foot
        raise ValueError("第 %d 页是深色页但认不出 kind" % page.idx)

    # ---- sources：正文起点比别人低 0.35 ----
    src = page.find(x=MARGIN, y=BODY_TOP + 0.35, w=13.333 - 2 * MARGIN)
    if src:
        d = {"kind": "sources", "title": title_box["text"], "groups": []}
        if lead:
            d["lead"] = lead["text"]
        cur = None
        for t, bold in zip(src["paras"], src["bolds"]):
            if not t:
                continue
            if bold:
                cur = {"h": t, "items": []}
                d["groups"].append(cur)
            elif cur is not None:
                cur["items"].append(t)
        return d, foot

    # ---- toc：标题是「导览」、三栏 ----
    if title_box["text"] == "导览":
        d = {"kind": "toc", "items": []}
        cols = sorted(page.all(y=BODY_TOP), key=lambda b: b["x"])
        for c in cols:
            paras = [(t, b) for t, b in zip(c["paras"], c["bolds"]) if t]
            i = 0
            while i < len(paras):
                t, bold = paras[i]
                if not bold:
                    i += 1
                    continue
                m = re.match(r"^(\S+)  (.*)$", t)
                it = {"no": m.group(1), "title": m.group(2)} if m else {"no": "", "title": t}
                if i + 1 < len(paras) and not paras[i + 1][1]:
                    it["gist"] = paras[i + 1][0]
                    i += 1
                d["items"].append(it)
                i += 1
        note = page.find(x=MARGIN, y=6.45)   # 表格页 12.09 宽、导览页 11.6 宽，只认位置
        if note:
            d["note"] = note["text"]
        return d, foot

    body = page.find(x=MARGIN, y=BODY_TOP, w=13.333 - 2 * MARGIN)
    if body is None:
        raise ValueError("第 %d 页没有正文区，认不出 kind" % page.idx)
    paras = [t for t in body["paras"] if t]

    # ---- goals / qa / summary：同一个正文框，靠段首标记分 ----
    if paras and paras[0].startswith("▸ "):
        d = base("goals")
        d["goals"] = [strip_bullet(t, "▸") for t in paras]
        return d, foot
    if paras and paras[0].startswith("Q  "):
        d = base("qa")
        d["qa"] = []
        for t in paras:
            if t.startswith("Q  "):
                d["qa"].append({"q": t[3:], "a": ""})
            elif t.startswith("A  ") and d["qa"]:
                d["qa"][-1]["a"] = t[3:]
        return d, foot
    if paras and re.match(r"^\d+　", paras[0]):
        d = base("summary")
        d["takeaways"] = [re.sub(r"^\d+　", "", t) for t in paras]
        nxt = [r for r in page.rects_fill(CARD) if r["text"].startswith("下一章　")]
        if nxt:
            d["next"] = nxt[0]["text"][len("下一章　"):]
        return d, foot
    raise ValueError("第 %d 页正文段首认不出（首段：%s）" % (page.idx, paras[:1]))


def extract(path):
    prs = Presentation(path)
    module = os.path.basename(path).replace("-讲义.pptx", "")
    spec = {"module": module, "pages": []}
    # 页脚页码是**整册**开关（kb_deck_build v9.0 起默认关）：任一页有 `p.NN` 就登记开。
    with_pn = any(
        any(sh.has_text_frame and re.match(r"^p\.\d+$", sh.text_frame.text.strip())
            for sh in sl.shapes)
        for sl in prs.slides)
    if with_pn:
        spec["page_numbers"] = True
    for i, slide in enumerate(prs.slides, 1):
        page, foot = parse(Page(slide, i), module)
        if "_display" in page:
            spec["display"] = page.pop("_display")
        if foot is not None:
            # footer() 里 foot 优先于 chapter_id；重渲染只需要 foot 就够。
            # 能拆出「模块 · 章节ID」形态的，额外留 chapter_id 供门禁与人读。
            if foot.startswith(module + " · "):
                page["chapter_id"] = foot[len(module) + 3:]
            else:
                page["foot"] = foot
        spec["pages"].append(page)
    return spec


def verify(src_pptx, spec):
    """往返自检：重渲染 → 逐页文本比对。原册字节不动。"""
    sys.path.insert(0, HERE)
    import kb_deck_build
    tmp = tempfile.mkdtemp(prefix="deckrt-")
    out = os.path.join(tmp, "rebuilt.pptx")
    kb_deck_build.build(spec, out)

    def dump(p):
        r = subprocess.run([sys.executable, os.path.join(HERE, "kb_deck_text.py"), p],
                           capture_output=True, text=True)
        return r.stdout

    a, b = dump(src_pptx), dump(out)
    if a == b:
        print("  往返一致：重渲染的 %d 页与原册逐页文本完全相同。" % len(spec["pages"]))
        return True
    la, lb = a.split("\n"), b.split("\n")
    print("  往返不一致（原册 %d 行 / 重渲染 %d 行），前几处差异：" % (len(la), len(lb)))
    shown = 0
    for i in range(max(len(la), len(lb))):
        x = la[i] if i < len(la) else "(无)"
        y = lb[i] if i < len(lb) else "(无)"
        if x != y:
            print("    行 %d\n      原: %s\n      新: %s" % (i + 1, x[:90], y[:90]))
            shown += 1
            if shown >= 6:
                break
    return False


def main(argv):
    if len(argv) < 3:
        print(__doc__)
        return 2
    src, dst = argv[1], argv[2]
    try:
        spec = extract(src)
    except ValueError as e:
        print("反解失败：%s" % e)
        return 1
    with open(dst, "w", encoding="utf-8") as f:
        json.dump(spec, f, ensure_ascii=False, indent=1)
        f.write("\n")
    kinds = {}
    for p in spec["pages"]:
        kinds[p["kind"]] = kinds.get(p["kind"], 0) + 1
    print("已反解 %s → %s：%d 页" % (src, dst, len(spec["pages"])))
    print("  kind 分布：%s" % "、".join("%s×%d" % kv for kv in sorted(kinds.items())))
    if "--verify" in argv:
        return 0 if verify(src, spec) else 1
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
